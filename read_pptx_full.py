from pptx import Presentation

prs = Presentation("OpenEnv rreq.pptx")
for i, slide in enumerate(prs.slides):
    print(f"\n{'='*80}")
    print(f"SLIDE {i+1}")
    print(f"{'='*80}")
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text.strip():
            print(shape.text)
