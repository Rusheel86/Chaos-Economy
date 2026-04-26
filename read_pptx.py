from pptx import Presentation

try:
    prs = Presentation("OpenEnv rreq.pptx")
    for i, slide in enumerate(prs.slides):
        print(f"--- Slide {i+1} ---")
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                print(shape.text)
except Exception as e:
    print("Error:", e)
